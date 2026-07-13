import os
import re
import base64
import io
import tempfile
import logging
from enum import Enum
import fitz
from PIL import Image
from ..config import settings

logger = logging.getLogger("visibility-docs")


class FileType(str, Enum):
    DIGITAL_PDF = "digital_pdf"
    SCANNED_PDF = "scanned_pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    TXT = "txt"
    IMAGE = "image"


ALLOWED_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"}


def _is_image_ext(ext: str) -> bool:
    return ext.lower() in ALLOWED_IMAGE_EXT


def detect_file_type(file_path: str) -> FileType:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".docx":
        return FileType.DOCX
    if ext == ".xlsx":
        return FileType.XLSX
    if ext == ".pptx":
        return FileType.PPTX
    if ext == ".txt":
        return FileType.TXT
    if _is_image_ext(ext):
        return FileType.IMAGE
    if ext == ".pdf":
        try:
            doc = fitz.open(file_path)
            page_count = doc.page_count
            text_len = 0
            for page in doc:
                text_len += len(page.get_text().strip())
            doc.close()
            avg_chars = text_len / max(page_count, 1)
            if avg_chars > 50:
                return FileType.DIGITAL_PDF
            return FileType.SCANNED_PDF
        except Exception as e:
            logger.warning(f"PDF detection failed for {file_path}: {e}")
            return FileType.SCANNED_PDF

    logger.warning(f"Unknown file type for {file_path}, treating as scanned")
    return FileType.SCANNED_PDF


def _extract_digital_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    pages = []
    for page_num, page in enumerate(doc, 1):
        blocks = page.get_text("dict").get("blocks", [])
        page_lines = []
        for block in blocks:
            btype = block.get("type", 0)
            if btype == 0:
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue
                    text = "".join(s.get("text", "") for s in spans)
                    font_size = max(s.get("size", 12) for s in spans)
                    flags = spans[0].get("flags", 0)
                    is_bold = bool(flags & 16) or "bold" in spans[0].get("font", "").lower()

                    if is_bold and font_size > 11:
                        text = f"## {text.strip()}"
                    page_lines.append(text.strip())
            elif btype == 1:
                page_lines.append("[IMAGE]")
        page_text = "\n".join(line for line in page_lines if line)
        pages.append(f"<!-- Page {page_num} -->\n{page_text}")
    doc.close()
    result = "\n\n".join(pages)
    result = _normalize_markdown(result)
    return result


def _extract_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        style = para.style.name.lower() if para.style else ""
        text = para.text.strip()
        if not text:
            parts.append("")
            continue
        if "heading" in style or "title" in style:
            try:
                level = int(re.search(r"heading\s*(\d+)", style).group(1))
                parts.append(f"{'#' * level} {text}")
            except Exception:
                parts.append(f"## {text}")
        elif any(r.bold for r in para.runs if r.bold):
            parts.append(f"**{text}**")
        else:
            parts.append(text)

    tables = []
    for table in doc.tables:
        md_rows = []
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            md_rows.append("| " + " | ".join(cells) + " |")
            if row_idx == 0:
                md_rows.append("| " + " | ".join("---" for _ in cells) + " |")
        tables.append("\n".join(md_rows))

    result = "\n".join(parts)
    if tables:
        result += "\n\n" + "\n\n".join(tables)
    result = _normalize_markdown(result)
    return result


def _extract_txt(file_path: str) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeError, Exception):
            continue
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows_data.append(cells)
        if rows_data:
            ncols = max(len(r) for r in rows_data)
            sep = "| " + " | ".join(["---"] * ncols) + " |"
            parts.append("| " + " | ".join(rows_data[0]) + " |")
            parts.append(sep)
            for cells in rows_data[1:]:
                parts.append("| " + " | ".join(cells) + " |")
        parts.append("")
    wb.close()
    result = "\n".join(parts)
    return _normalize_markdown(result)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Slide {slide_idx} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                table = shape.table
                md_rows = []
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    md_rows.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                slide_texts.append("\n".join(md_rows))
        parts.append("\n".join(slide_texts))
    result = "\n\n".join(parts)
    return _normalize_markdown(result)


def _page_to_image(page) -> str:
    pix = page.get_pixmap(dpi=150)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    target_w = 768
    w_percent = target_w / float(img.width)
    h_size = int(float(img.height) * float(w_percent))
    img = img.resize((target_w, h_size), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=80, optimize=True)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def _load_image_b64(file_path: str) -> str:
    target_w = 768
    img = Image.open(file_path)
    if img.mode != "RGB":
        img = img.convert("RGB")
    w_percent = target_w / float(img.width)
    h_size = int(float(img.height) * float(w_percent))
    img = img.resize((target_w, h_size), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=80, optimize=True)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


VISION_PROMPT = """You are an OCR engine. Extract ALL visible text from this document page image.

Rules:
- Extract every character, word, number exactly as it appears
- Preserve layout: headings, paragraphs, bullet lists, numbered lists, tables
- Output tables as Markdown tables
- If the image contains a diagram, engineering drawing, chart, or figure, extract all visible labels and annotations
- Prefix each page with '--- Page X ---'
- Output ONLY clean Markdown — no explanations, no summaries, no commentary
- NEVER convert tables into paragraphs — always use Markdown table syntax
- Preserve figure labels, captions, warnings, notes, cautions as they appear"""


def _vision_ocr(image_b64s: list[str]) -> str:
    from .groq_service import groq_service
    if not groq_service.available or not groq_service.vision_available:
        logger.warning("Groq vision not available")
        return ""

    texts = []
    batch_size = 5
    for i in range(0, len(image_b64s), batch_size):
        batch = image_b64s[i:i + batch_size]
        content = [{"type": "text", "text": VISION_PROMPT}]
        for b64 in batch:
            content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}})
        try:
            result = groq_service.chat_vision(
                [{"role": "user", "content": content}],
                temperature=0.0,
                max_tokens=8192,
            )
            if result and not result.startswith("[Groq"):
                texts.append(result)
            else:
                texts.append("")
        except Exception as e:
            logger.warning(f"Vision batch {i//batch_size + 1} failed: {e}")
            texts.append("")

    return "\n\n".join(texts)


def process_scanned_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    b64s = []
    for page in doc:
        b64s.append(_page_to_image(page))
    doc.close()

    if not b64s:
        return ""

    text = _vision_ocr(b64s)
    if text.strip():
        text = _normalize_markdown(text)
        return text

    logger.warning("Vision returned empty for scanned PDF, trying PyMuPDF fallback")
    try:
        doc = fitz.open(file_path)
        parts = []
        for page in doc:
            t = page.get_text().strip()
            if t:
                parts.append(t)
        doc.close()
        fallback = "\n\n".join(parts)
        return _normalize_markdown(fallback) if fallback else "[OCR failed]"
    except Exception:
        return "[OCR failed]"


def process_image(file_path: str) -> str:
    b64 = _load_image_b64(file_path)
    text = _vision_ocr([b64])
    if text.strip():
        return _normalize_markdown(text)
    return "[OCR failed]"


def _normalize_markdown(text: str) -> str:
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\n{4,}', '\n\n\n', text)
    text = re.sub(r'[ \t]+$', '', text, flags=re.MULTILINE)
    lines = text.split('\n')
    result = []
    for line in lines:
        stripped = line.strip()
        if re.match(r'^#{1,6}\s', stripped):
            result.append(stripped)
        elif re.match(r'^[-*]\s', stripped):
            result.append(stripped)
        elif re.match(r'^\d+[.)]\s', stripped):
            result.append(stripped)
        elif re.match(r'^\|', stripped):
            result.append(stripped)
        elif re.match(r'^>\s', stripped):
            result.append(stripped)
        elif re.match(r'^```', stripped):
            result.append(stripped)
        else:
            result.append(line)
    return '\n'.join(result)


def process_document(file_path: str) -> dict:
    file_type = detect_file_type(file_path)
    logger.info(f"[OCR] Detected: {file_type.value} — {file_path}")

    text = ""
    page_count = 0

    if file_type == FileType.DIGITAL_PDF:
        text = _extract_digital_pdf(file_path)
        if text:
            doc = fitz.open(file_path)
            page_count = doc.page_count
            doc.close()
        source = "digital_pdf"

    elif file_type == FileType.DOCX:
        text = _extract_docx(file_path)
        page_count = max(1, len(text) // 2000)
        source = "docx"

    elif file_type == FileType.XLSX:
        text = _extract_xlsx(file_path)
        page_count = max(1, len(text) // 2000)
        source = "xlsx"

    elif file_type == FileType.PPTX:
        text = _extract_pptx(file_path)
        page_count = max(1, len(text) // 2000)
        source = "pptx"

    elif file_type == FileType.TXT:
        text = _extract_txt(file_path)
        page_count = max(1, len(text) // 2000)
        source = "txt"

    elif file_type == FileType.SCANNED_PDF:
        text = process_scanned_pdf(file_path)
        if text and text != "[OCR failed]":
            try:
                doc = fitz.open(file_path)
                page_count = doc.page_count
                doc.close()
            except Exception:
                page_count = 0
        source = "vision"

    elif file_type == FileType.IMAGE:
        text = process_image(file_path)
        page_count = 1
        source = "vision"

    logger.info(f"[OCR] Result: {len(text)} chars, {page_count} pages, source={source}")
    return {"text": text, "page_count": page_count, "source": source}


ocr_service = type("OCRService", (), {"process_document": staticmethod(process_document)})()
